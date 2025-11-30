import os
from pptx import Presentation

def extract_content_from_ppt(file_path: str):
    """
    解析PPT，返回每一页的文本内容列表
    格式: [{"page": 1, "text": "xxx"}, ...]
    """
    prs = Presentation(file_path)
    slides_content = []

    for i, slide in enumerate(prs.slides):
        page_text = []
        
        # --- 修复开始：安全获取标题 ---
        title_text = ""
        # 只有当标题存在，且有文本框时，才获取标题
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title_text = slide.shapes.title.text.strip()
            page_text.append(title_text)
        # --- 修复结束 ---
        
        # 提取正文和其他文本框
        for shape in slide.shapes:
            # 1. 检查形状是否有文本框属性 (has_text_frame)
            if not shape.has_text_frame:
                continue
            
            # 2. 获取文本
            text = shape.text.strip()
            
            # 3. 过滤空文本
            if not text:
                continue

            # 4. 避免重复：如果这段字和标题一样，就不要再加一遍了
            # 注意：这里我们比较的是变量 title_text，而不是直接调 slide.shapes.title.text
            if title_text and text == title_text:
                continue
                
            page_text.append(text)
        
        full_text = "\n".join(page_text).strip()
        
        # 只有这一页确实解析出了文字才返回
        if full_text:
            slides_content.append({
                "page": i + 1,
                "text": full_text
            })
            
    return slides_content